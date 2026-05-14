#!/usr/bin/env python3
"""update-modernization-docs.py — Apply 2026-05-14 progress update to:
  - docs/HairMNL-90-Day-Modernization-Roadmap.docx
  - docs/HairMNL-Theme-Modernization-Deck.pptx

Inputs: the existing docs (last updated 2026-05-06 — commit b72f583).
Output: updates in place. Run `git diff --stat docs/` to see what changed.

What this script changes (incremental — preserves all existing content):
  Roadmap (.docx):
    1. Title date: "Updated May 6, 2026" → "Updated May 14, 2026"
    2. Phase 1 heading: "Apr 26 → May 6" → "Apr 26 → May 14"
    3. Inserts new Heading 2 sections AFTER existing Phase 1 sub-sections,
       documenting May 7 → May 14 shipped work:
         - CLS Phase 2 sprint (mobile + desktop)
         - kt0 CSS regression prevention (rule, lint, smoke test)
         - GTM/RUM observability hardening (d00 metric_rating, orphan triggers)
         - Third-party app rationalization (BSS Lock+Solution+LTAP, BTA, VQB)
         - Dashboard refinements (Omura desktop charts, 3-band zones, noise band)
    4. Adds an "At-a-glance" update note with new CrUX numbers
    5. Updates decision statuses (3q4 still open; ijh + ebh closed; BSS L+S+LTAP done)
    6. Updates app stack table to mark removed apps
    7. Updates Forward Plan dates

  Deck (.pptx):
    1. Cover slide date: May 6 → May 14
    2. At-a-glance slide 2 numbers refreshed
    3. Adds 4 new slides after the existing Phase 1 slides covering the May 7-14 work
    4. Closing summary slide 26: updates decision status list

The script is idempotent at the date-line level (won't double-update) but the
section-insertion logic is NOT — running twice would duplicate the inserted
sections. After running once, the doc is at "May 14" state.
"""

import os
import sys
from datetime import date
from pathlib import Path

import docx
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
from copy import deepcopy

REPO = Path(__file__).resolve().parent.parent
ROADMAP = REPO / "docs" / "HairMNL-90-Day-Modernization-Roadmap.docx"
DECK = REPO / "docs" / "HairMNL-Theme-Modernization-Deck.pptx"


# =============================================================================
# Content blocks (what to insert) — single source of truth, used by both docs.
# =============================================================================

NEW_SECTIONS_ROADMAP = [
    {
        "heading": "What we shipped — May 7 → May 14 (Phase 2 sprint)",
        "intro": (
            "Eight days, ~50 Beads issues closed, 11 substantive commits. The biggest "
            "field-data win is on desktop CLS: real-shopper CrUX p75 good-rate climbed "
            "from 30.5% to 56.3% (+25.8 pp) — half of remaining desktop shoppers are now "
            "in the 'good' bucket vs. less than a third 8 days ago. Mobile CLS p75 also "
            "improved (74.5% → 82.3%, +7.8 pp)."
        ),
        "bullets": [],
    },
    {
        "heading": "Phase 2 CLS sprint — desktop focus",
        "intro": "Six Phase 2 fixes shipped specifically targeting the desktop-only contributors that the May 2 mobile sprint did not touch.",
        "bullets": [
            "jvf — Cart-page quantity + price cell min-heights (resolves cart-row CLS on /cart desktop).",
            "p52 — Judge.me PDP review block + collection-page preview badge min-heights (220 px + 18 px reservations).",
            "dzi (then hod) — BOGOS free-gift cart container reservation: 80 px → 160 px → 240 px. The May 13 hod update covers heavy-promo carts with 3 gift rows; before the bump, those carts shifted by ~220 px.",
            "38l — LimeSpot recommendation card image aspect-ratio + object-fit lock.",
            "ojx — Pro-blogger related-articles thumbnail dimensions added (was 136 URLs in the desktop poor bucket).",
            "z9y — PDP Judge.me preview badge reservation + gallery zoom yield (covered 884 URLs).",
            "xn4 — Article-cards Option-B template-body scope (replaces an earlier global fix that risked over-applying).",
            "987 Phase 2 #1 — Blog article body iframe/video aspect-ratio reservation (16/9 default, covers YouTube/Vimeo/Instagram embeds).",
            "hod — Cart-template (cart__te) 84-event desktop CLS contributor: BOGOS reservation increased to 240 px (the actual winning rule lives in custom-theme.css, also updated).",
            "7i0 — Announcement bar ticker height lock — pin inner [data-ticker-frame] to 24 px + overflow:hidden + white-space:nowrap to eliminate the font-FOUT inner shift.",
        ],
    },
    {
        "heading": "kt0 CSS containment rule — locked in as policy + tooling",
        "intro": (
            "On May 11 a kt0-class regression broke the cart drawer (commit ff1ef80) and "
            "then on May 12 broke the Reamaze chat widget — both caused by adding contain:layout "
            "to a parent of a fixed-position descendant. Three layers of prevention are now in place."
        ),
        "bullets": [
            "Rule documented in CLAUDE.md + .opencode_hints: containing-block-creating properties (contain:layout/paint, transform, filter, backdrop-filter, perspective, will-change:transform) on a parent of any fixed/absolute descendant silently break the overlay's positioning math.",
            "scripts/check-overlay-css.py — Python lint that scans all .liquid + .css + .scss for the bug pattern; runs in CI on every PR touching those files via .github/workflows/kt0-css-lint.yml.",
            "scripts/smoke-test-drawers.py — Playwright-based test that programmatically opens the cart drawer, mobile nav, search overlay, and modals on the draft theme and walks each fixed-position parent chain looking for the bug pattern. Required to pass before any live push touching layout-sensitive CSS.",
            "Process change: any agent (Claude Code, OpenCode, Cursor) editing CSS reservations must grep BOTH snippets/css-overrides.liquid AND assets/custom-theme.css for the selector — duplicate rules with different specificities caused the May 13 hod fix to ship as a NO-OP until the coordinator's computed-style smoke test caught it. The 'CSS reservation parity' hint codifies this.",
        ],
    },
    {
        "heading": "GTM + RUM observability hardening",
        "intro": (
            "Three observability bugs were resolved that had been masking real-shopper "
            "Core Web Vitals data — about half of every metric event was previously "
            "arriving in GA4 without its metric_rating dimension populated."
        ),
        "bullets": [
            "d00 — Found and fixed a gtag-double-fire pattern in snippets/web-vitals-reporter.liquid. When GTM is loaded, window.gtag('event', 'LCP', params) was sending the metric BOTH directly to GA4 (correct rating) AND via GTM's internal gtag bridge to tag 766 (where the DLV reads returned empty string because the bridge doesn't populate the flat dataLayer model). 24h after the fix landed, the 'unknown' metric_rating bucket dropped from 50–60% to 2–6% across all five vitals. GA4 RUM now reports clean ratings for the first time.",
            "vux — Diagnosed and reverted a GTM v141 regression that had been spiking the LCP unknown bucket on the dashboard.",
            "t4m, 4r8, miy — Three GTM tags that referenced an orphan trigger ID (2147479553) were repaired/disabled (Klaviyo Form Submission Listener, Conversion Linker, TrafficGuard).",
            "p8j — JS error reporter enhanced to capture e.error.name for richer error_type categorization (replaces the catch-all 'error' bucket that was 39k events/7d).",
        ],
    },
    {
        "heading": "Third-party app rationalization — net code reduction",
        "intro": (
            "Three vendor apps confirmed uninstalled on Shopify; dead theme code removed; "
            "settings_data.json + caller blocks cleaned. Net effect: 49 dead snippet files "
            "deleted, every PDP no longer ships two dead fetch() calls, BTA bootstrap fully "
            "blocked. Approximately 2,021 lines of dead Liquid removed in a single commit pair."
        ),
        "bullets": [
            "BSS B2B Lock (uninstalled) — 13 bsscommerce-* / bss-passcode-* / bss-custom-login snippets removed (bd dyp).",
            "BSS B2B Solution (uninstalled) — 20 bss-b2b-* snippets removed (covers VAT pricing, tax-cart, wholesaler register, state config) + the bss-b2b-wholesale-solution app block reference (bd dyp).",
            "BSS LTAP / Login To Access Page (uninstalled) — 16 bss-ltap-* + bsscommerce-redirect-product-page-logic snippets removed + 2 duplicate inline scripts at the top of every PDP that POSTed to login-to-access-page-api.bsscommerce.com (bd gj0).",
            "BookThatApp (uninstalled) — 60-line createElement guard removed from layout/theme.liquid + BTA app block reference removed from settings_data.json (bd p84).",
            "VisualQuizBuilder (uninstalled) — same createElement guard pattern removed; was previously gated to /pages/*-quiz routes (bd 8z9).",
            "Customer Portal preserved — bss-b2b-cp-* + shared bss-b2b-js/jquery-341-js + currency-format + featured-product-vat-styles snippets stay, gated behind 'content_for_header contains bss-b2b-cp' + customer.tags contains 'backbar-approved' (Decision 4 still pending founder call — bd bjk).",
        ],
    },
    {
        "heading": "Dashboard refinements (Omura backport)",
        "intro": "Backported four dashboard improvements from the Omura theme to make trend reading easier.",
        "bullets": [
            "z3x epic — Desktop trend charts added + mobile/desktop tab toggle. Previously the trend grid only showed mobile.",
            "bk2 — Replaced threshold lines with full 3-band Web Vitals zones (good/needs-improvement/poor) on every CWV trend chart.",
            "c0w — PSI lab cards now show daily median + a min/max noise band instead of a single noisy run. PSI single-run is famously variable for HairMNL (last 5 runs spanned score 28–40); the median+band stops one bad run from looking like a regression.",
            "u8w — GHA dashboard refresh bumped from 2 to 3 PSI runs per snapshot.",
            "Critical-CSS audit (ax0) ran across 6 templates against the full 344 KB bundle — confirmed 39 KB inline is correctly sized.",
        ],
    },
    {
        "heading": "Critical hotfixes (off the Phase 2 critical path but shipped this window)",
        "intro": "Three user-reported live regressions surfaced and were resolved within hours each.",
        "bullets": [
            "lki (2026-05-12) — Reamaze chat widget broken by an over-broad CSS containment rule (7fz) that applied contain:layout to a parent of fixed-position Reamaze descendants. Narrowed the rule to #reamaze-widget-label only; chat widget restored.",
            "fhh (2026-05-13) — Mobile /cart layout broken: assets/cart-page.css extracted from theme.css on 2026-04-26 had its two @media wrappers stripped. The desktop 5-column .cart__items__grid was applying at every viewport, forcing 5 columns at 390 px (image collapsed, titles wrapped 6+ lines). Restored both media-query wrappers; mobile cart now uses the 2-column stacked layout, desktop preserves the 5-column grid.",
            "Cart drawer cross-page rendering bug (May 11) — preceded the kt0 codification work above.",
        ],
    },
]


# =============================================================================
# Roadmap (.docx) update
# =============================================================================

def update_roadmap():
    doc = docx.Document(str(ROADMAP))

    # --- 1. Update title date line ---
    for p in doc.paragraphs:
        if p.text.strip() == "Updated May 6, 2026":
            for run in p.runs:
                run.text = ""
            p.runs[0].text = f"Updated May 14, 2026"
            break

    # --- 2. Update Phase 1 heading date range ---
    for p in doc.paragraphs:
        if "Phase 1 — what's already shipped" in p.text:
            for run in p.runs:
                if "May 6" in run.text:
                    run.text = run.text.replace("May 6", "May 14")
            # If the date didn't appear in a single run, replace the whole text
            if "May 6" in p.text:
                full = p.text.replace("May 6", "May 14")
                # Clear and rewrite (preserves the first run's style)
                for r in p.runs[1:]:
                    r.text = ""
                p.runs[0].text = full
            break

    # --- 3. Insert new Phase 2 sections AFTER the last Phase 1 sub-section ---
    # Find the "Decisions pending" heading — that's where we want to insert just BEFORE.
    decisions_idx = None
    for i, p in enumerate(doc.paragraphs):
        if (p.style.name if p.style else "") == "Heading 1" and "Decisions pending" in p.text:
            decisions_idx = i
            break
    if decisions_idx is None:
        raise RuntimeError("Couldn't locate 'Decisions pending' heading")

    # python-docx doesn't have a clean 'insert paragraph before X' for body
    # elements; we have to manipulate the XML directly.
    from docx.oxml.ns import qn
    target_p = doc.paragraphs[decisions_idx]
    target_elem = target_p._p
    body = target_elem.getparent()

    # Find template paragraphs we can clone for each style (avoids the
    # python-docx style-resolution bug on this doc)
    template_para = {}
    for p in doc.paragraphs:
        style_name = p.style.name if p.style else "Normal"
        if style_name in ("Heading 2", "Normal", "List Paragraph") and style_name not in template_para:
            template_para[style_name] = p

    def new_para_before(target, text, style="Normal"):
        # Clone a template paragraph of the same style, clear its runs, insert text
        template = template_para.get(style) or template_para.get("Normal")
        if template is None:
            # Fallback — use add_paragraph (may fail with style cache bug)
            p = doc.add_paragraph(text)
            new_elem = p._p
            body.remove(new_elem)
            target.addprevious(new_elem)
            return p
        new_elem = deepcopy(template._p)
        # Strip all existing run elements but keep paragraph properties (style)
        from docx.oxml.ns import qn
        for r in list(new_elem.findall(qn("w:r"))):
            new_elem.remove(r)
        # Add a single fresh run with the new text
        run_xml = template._p.findall(qn("w:r"))
        if run_xml:
            new_run = deepcopy(run_xml[0])
            # Clear text in cloned run, set new
            for t in new_run.findall(qn("w:t")):
                new_run.remove(t)
            from lxml import etree
            t_elem = etree.SubElement(new_run, qn("w:t"))
            t_elem.text = text
            t_elem.set(qn("xml:space"), "preserve")
            new_elem.append(new_run)
        target.addprevious(new_elem)
        # Return a Paragraph wrapper
        from docx.text.paragraph import Paragraph
        return Paragraph(new_elem, doc.paragraphs[0]._parent)

    # Sentinel — only insert sections if not already present (idempotency check)
    already_inserted = any(
        "Phase 2 sprint" in p.text and (p.style.name if p.style else "").startswith("Heading")
        for p in doc.paragraphs
    )
    if not already_inserted:
        for section in NEW_SECTIONS_ROADMAP:
            new_para_before(target_elem, section["heading"], style="Heading 2")
            if section.get("intro"):
                new_para_before(target_elem, section["intro"], style="Normal")
            for b in section.get("bullets", []):
                new_para_before(target_elem, b, style="List Paragraph")

    # --- 4. Update App stack table (Table 2) to mark recently-removed apps ---
    if len(doc.tables) >= 3:
        t = doc.tables[2]
        for row in t.rows[1:]:  # skip header
            cells = row.cells
            app = cells[0].text.strip()
            status = cells[2].text.strip() if len(cells) > 2 else ""
            action = cells[3].text.strip() if len(cells) > 3 else ""
            if "Subscribe It" in app or "Treedify" in app or "Hextom Conversion" in app:
                cells[2].text = "Uninstalled May 9"
                cells[3].text = "Done. Dead theme code removed (bd ijh)."
            elif "BSS B2B" in app and "Customer Portal" in app:
                cells[2].text = "Founder decision pending [bjk]"
                cells[3].text = "JS gated to pro logins only; zero perf cost on guest pageloads. Lock + Solution + LTAP variants of BSS uninstalled May 13 (49 dead snippets removed, bd dyp + gj0)."

    # --- 5. Update Forward Plan table (Table 1) — dates are stale ---
    if len(doc.tables) >= 2:
        t = doc.tables[1]
        new_rows = [
            ("May 13 → 19", "Phase 2 CLS desktop fixes (hod, 7i0 — cart-template + announcement bar shipped May 13). Watch GA4 RUM 7-day window clear of pre-d00-fix data by May 19.", "Web team"),
            ("May 14 → 20", "Continue Phase 2 CLS dispatches via OC sub-orchestrator (uj2 #MainContent umbrella; r8r STKY replacement; 15f Elevar audit). LimeSpot tweaks (Decision 1) — pending.", "Web team + Marketing"),
            ("May 20 → 26", "Vertex Phase 1 — catalog sync (Cloud Functions). Re-verify mobile CLS sustained <0.10 in CrUX. CrUX desktop CLS first read at trailing edge of fixes.", "Web team (Vertex separate session)"),
            ("May 27 → Jun 2", "Vertex Phase 2 — event tracking pipeline. Phase 2 CLS verification: cart-template event count drop to <20/7d.", "Web team"),
            ("Jun 3 → Jun 10", "Vertex Phase 3 — cutover with 14-day monitoring. 987 desktop CLS final read (28-day CrUX window slides past 2026-05-13 commits ~ Jun 10).", "Founders + Web team"),
        ]
        # Replace existing data rows (keep header)
        # Add rows if needed
        existing_rows = len(t.rows) - 1
        for i, (week, work, owner) in enumerate(new_rows):
            if i + 1 < len(t.rows):
                row = t.rows[i + 1]
            else:
                row = t.add_row()
            row.cells[0].text = week
            row.cells[1].text = work
            row.cells[2].text = owner

    # --- 6. Update the Decisions section: mark Decision 3 sub-items closed ---
    # Decision 3 sub-bullets — find and update statuses
    for p in doc.paragraphs:
        txt = p.text.strip()
        if "Judge.me admin (issue ebh)" in txt:
            # Mark closed
            if "DONE" not in txt:
                p.runs[0].text = "Judge.me admin (issue ebh) — DONE (closed 2026-05-11). Lite mode enabled, reviews-per-page reduced, unused widgets disabled. Live impact expected at next CrUX read."
                for r in p.runs[1:]:
                    r.text = ""
        elif "Verify Subscribe It / Treedify / Hextom CTB install status (issue ijh)" in txt:
            if "DONE" not in txt:
                p.runs[0].text = "Verify Subscribe It / Treedify / Hextom CTB install status (issue ijh) — DONE (closed 2026-05-09). All three confirmed uninstalled; dead theme code removed (commit 0eb4f59)."
                for r in p.runs[1:]:
                    r.text = ""

    # Decision 4 update — Lock + Solution + LTAP are done; only Customer Portal left
    for p in doc.paragraphs:
        if "Status: PENDING [bjk] — founder decision required" in p.text:
            new_text = (
                "Status: PENDING [bjk] — founder decision required. The non-Customer-Portal "
                "BSS apps (B2B Lock, B2B Solution, LTAP) were uninstalled on Shopify and dead "
                "theme code removed on 2026-05-13 (49 snippet files deleted, bd dyp + gj0). "
                "Customer Portal is the only remaining BSS app — JS already gated to backbar-"
                "approved customers; zero perf cost on guest pageloads. The decision below is "
                "now narrower: keep Customer Portal vs. uninstall just it."
            )
            for r in p.runs[1:]:
                r.text = ""
            p.runs[0].text = new_text
            break

    # --- 7. Append a brief "At-a-glance — May 14 refresh" paragraph after Table 0 ---
    # Table 0 is the AT A GLANCE table. We'll add a paragraph right after it.
    if len(doc.tables) >= 1:
        t0 = doc.tables[0]
        # Update the existing row with a fresh CrUX snapshot line
        if len(t0.rows) >= 2:
            existing = t0.rows[1].cells[0].text
            if "May 14 update" not in existing:
                t0.rows[1].cells[0].text = existing + (
                    "\n\nMay 14 update — CrUX p75 real-shopper field data: mobile LCP good 61.9% / "
                    "CLS good 82.3% / INP good 67.9%. Desktop LCP good 72.8% / CLS good 56.3% / "
                    "INP good 84.2%. Desktop CLS climbed +25.8 pp since May 6 (30.5% → 56.3% good) — "
                    "the biggest Phase 2 win. GA4 metric_rating data is now clean (~5% unknown vs "
                    "~50% pre-d00-fix on May 13) so future Phase 2 dispatches have reliable RUM signal."
                )

    doc.save(str(ROADMAP))
    print(f"Saved {ROADMAP.relative_to(REPO)}")


# =============================================================================
# Deck (.pptx) update
# =============================================================================

def update_deck():
    pres = Presentation(str(DECK))

    # --- 1. Update cover slide date (slide 1) ---
    slide1 = pres.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "May 6, 2026" in run.text:
                        run.text = run.text.replace("May 6, 2026", "May 14, 2026")

    # --- 2. Refresh slide 2 (at-a-glance) numbers ---
    # Find slide 2 — has "Where we are — one-page read"
    for slide in pres.slides:
        head = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                head = shape.text_frame.text
                break
        if "Where we are" in head:
            # Update specific runs
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                tf = shape.text_frame
                for para in tf.paragraphs:
                    txt = para.text
                    if "−65%" in txt:
                        # update with May 14 mobile CLS good improvement headline
                        new_txt = "+25.8 pp"
                        if para.runs:
                            para.runs[0].text = new_txt
                            for r in para.runs[1:]:
                                r.text = ""
                    elif "−86%" in txt:
                        new_txt = "+7.8 pp"
                        if para.runs:
                            para.runs[0].text = new_txt
                            for r in para.runs[1:]:
                                r.text = ""
                    elif "vs Apr-26 baseline · TBT median 774 ms desktop" in txt:
                        new_txt = "Desktop CLS p75 good 30.5% → 56.3% (May 6 → May 14)"
                        if para.runs:
                            para.runs[0].text = new_txt
                            for r in para.runs[1:]:
                                r.text = ""
                    elif "TBT 5,600 ms → 774 ms" in txt:
                        new_txt = "Mobile CLS p75 good 74.5% → 82.3% (May 6 → May 14)"
                        if para.runs:
                            para.runs[0].text = new_txt
                            for r in para.runs[1:]:
                                r.text = ""
            break

    # --- 3. Insert 4 new slides after slide 10 (last "what we shipped" slide) ---
    # Slide layout 1 is typically the content layout
    layout = pres.slide_layouts[1] if len(pres.slide_layouts) > 1 else pres.slide_layouts[0]

    new_slides = [
        {
            "title": "What we shipped — May 7-14 sprint",
            "subtitle": "Phase 2 — 8 days, ~50 issues closed, 11 substantive commits",
            "bullets": [
                "Desktop CLS p75 good rate 30.5% → 56.3% (+25.8 pp) — biggest field win",
                "Mobile CLS p75 good rate 74.5% → 82.3% (+7.8 pp)",
                "GA4 metric_rating fix (d00) — unknown bucket 50% → 5% in 24h, RUM signal now clean",
                "49 dead snippet files deleted (BSS Lock + Solution + LTAP)",
                "kt0 CSS containment rule codified (lint + Playwright smoke test + .opencode_hints)",
                "Dashboard: desktop trend charts + 3-band CWV zones + PSI noise band",
            ],
        },
        {
            "title": "CLS Phase 2 — desktop focus",
            "subtitle": "Six fixes targeting desktop-only contributors the May 2 sprint did not touch",
            "bullets": [
                "jvf — Cart-page quantity + price cell min-heights",
                "p52 — Judge.me PDP review block (220 px) + collection preview badge (18 px)",
                "dzi → hod — BOGOS free-gift container 80 → 160 → 240 px (heavy-promo carts)",
                "38l — LimeSpot card image aspect-ratio + object-fit lock",
                "ojx — Pro-blogger related-articles thumbnail dimensions (136 URLs)",
                "z9y — PDP Judge.me preview badge + gallery zoom yield (884 URLs)",
                "987 #1 — Blog article body iframe/video aspect-ratio (YouTube/Vimeo/Instagram)",
                "7i0 — Announcement bar ticker height lock (font-FOUT inner shift)",
            ],
        },
        {
            "title": "kt0 CSS containment rule — three layers of prevention",
            "subtitle": "After two cart-drawer + Reamaze regressions in three days, this is now policy",
            "bullets": [
                "Rule: contain:layout / transform / filter / etc. on a parent of any fixed/abs descendant breaks overlay positioning silently",
                "Layer 1 — Lint: scripts/check-overlay-css.py runs in CI on every PR touching .liquid/.css/.scss",
                "Layer 2 — Smoke: scripts/smoke-test-drawers.py (Playwright) verifies cart drawer + mobile nav + search + modals on draft before live push",
                "Layer 3 — Process: CSS reservation parity hint — grep both css-overrides.liquid AND custom-theme.css before editing",
                "Caught the May 13 hod fix shipping as a NO-OP (selector-specificity trap) — computed-style smoke test required for every CSS reservation now",
            ],
        },
        {
            "title": "GTM observability + third-party app rationalization",
            "subtitle": "Real-shopper RUM data is now clean for the first time",
            "bullets": [
                "d00 — gtag double-fire fix: metric_rating unknown 50-60% → 2-6% in 24h (LCP/CLS/INP/FCP/TTFB all clean)",
                "vux — Reverted a GTM v141 instrumentation regression",
                "t4m + 4r8 + miy — Three GTM tags with orphan trigger IDs repaired (Klaviyo, Conversion Linker, TrafficGuard)",
                "BSS B2B Lock + Solution + LTAP uninstalled — 49 dead snippets removed, 2 dead fetch() calls per PDP eliminated",
                "BookThatApp + VisualQuizBuilder uninstalled — 60-line createElement guard + 2 app block refs removed",
                "Customer Portal preserved — gated to backbar-approved customers, zero guest perf cost (Decision 4 narrower now)",
            ],
        },
    ]

    # Find slide index after slide 10 (the code-health pass slide) to insert after
    # We add to the end and then reorder via XML manipulation
    from copy import deepcopy

    # Clone the structure of an existing "What we shipped" slide (slide 5 has
    # title + subtitle + body bullets in named shapes). The deck's CONTENT_MASTER
    # layout doesn't expose title/body placeholders, so we copy individual shape
    # elements from a real slide.
    from lxml import etree
    from pptx.oxml.ns import qn
    template_slide = pres.slides[4]  # slide 5 "What we shipped — JavaScript reduction"
    template_spTree = template_slide.shapes._spTree
    # Cache the individual <p:sp> elements from the template
    template_sp_elems = template_spTree.findall(qn("p:sp"))

    def _set_shape_text(shape, text):
        """Update a shape's text by directly rewriting its text frame XML.

        text_frame.text = "..." can fail silently on cloned slides because the
        deepcopy preserves run-level rPr formatting that python-pptx struggles
        to clear. Direct XML rewrite is reliable."""
        from docx.oxml.ns import qn as docx_qn  # noqa: F401
        from pptx.oxml.ns import qn
        tf = shape.text_frame
        txBody = tf._txBody
        # Remove all <a:p> children
        for p in list(txBody.findall(qn("a:p"))):
            txBody.remove(p)
        # Add fresh <a:p> with one <a:r> containing the text
        a_p = etree.SubElement(txBody, qn("a:p"))
        a_r = etree.SubElement(a_p, qn("a:r"))
        a_t = etree.SubElement(a_r, qn("a:t"))
        a_t.text = text

    def _set_body_text(shape, bullets):
        from pptx.oxml.ns import qn
        tf = shape.text_frame
        txBody = tf._txBody
        for p in list(txBody.findall(qn("a:p"))):
            txBody.remove(p)
        for b in bullets:
            a_p = etree.SubElement(txBody, qn("a:p"))
            a_r = etree.SubElement(a_p, qn("a:r"))
            a_t = etree.SubElement(a_r, qn("a:t"))
            a_t.text = b

    def _replace_text_in_sp(sp_elem, new_text):
        """Replace all text in a <p:sp> shape element. Operates at the XML level."""
        txBody = sp_elem.find(qn("p:txBody"))
        if txBody is None:
            return
        # Remove all existing <a:p> children
        for p in list(txBody.findall(qn("a:p"))):
            txBody.remove(p)
        # Add a fresh <a:p><a:r><a:t>new_text</a:t></a:r></a:p>
        a_p = etree.SubElement(txBody, qn("a:p"))
        a_r = etree.SubElement(a_p, qn("a:r"))
        a_t = etree.SubElement(a_r, qn("a:t"))
        a_t.text = new_text

    def _replace_body_bullets(sp_elem, bullets):
        """Replace bullets in a <p:sp> shape — one <a:p> per bullet."""
        txBody = sp_elem.find(qn("p:txBody"))
        if txBody is None:
            return
        for p in list(txBody.findall(qn("a:p"))):
            txBody.remove(p)
        for b in bullets:
            a_p = etree.SubElement(txBody, qn("a:p"))
            a_r = etree.SubElement(a_p, qn("a:r"))
            a_t = etree.SubElement(a_r, qn("a:t"))
            a_t.text = b

    def add_text_slide(title, subtitle, bullets):
        slide = pres.slides.add_slide(template_slide.slide_layout)
        target_spTree = slide.shapes._spTree
        # Clone shapes 0-3 from the template (title, subtitle, decorative-empty,
        # bullets). Skip shape 4 (the slide number placeholder which would
        # otherwise carry the literal "5" into every new slide).
        for idx, sp in enumerate(template_sp_elems[:4]):
            new_sp = deepcopy(sp)
            if idx == 0:
                _replace_text_in_sp(new_sp, title)
            elif idx == 1:
                _replace_text_in_sp(new_sp, subtitle)
            elif idx == 2:
                _replace_text_in_sp(new_sp, "")
            elif idx == 3:
                _replace_body_bullets(new_sp, bullets)
            target_spTree.append(new_sp)
        return slide

    # Sentinel — only insert if not already present
    titles = []
    for s in pres.slides:
        for shape in s.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if t:
                    titles.append(t)
                    break
    already_inserted = any("May 7-14 sprint" in t for t in titles)

    if not already_inserted:
        for spec in new_slides:
            add_text_slide(spec["title"], spec["subtitle"], spec["bullets"])

        # Reorder: move the 4 new slides to just after slide 10 (index 10 = slide 11 1-based)
        # The new slides are at the END now. We want them at indices 10, 11, 12, 13
        # (so that slide 11+ becomes the "What we shipped May 7-14" section)
        # python-pptx doesn't have direct reorder; manipulate sldIdLst directly.
        xml_slides = pres.slides._sldIdLst
        slides_list = list(xml_slides)
        # The 4 new slides are at the end (last 4 entries)
        target_insert_position = 10  # after slide 10 (0-based index 10 means slide 11 1-based)
        new_entries = slides_list[-4:]
        for entry in new_entries:
            xml_slides.remove(entry)
        # Re-insert in order at position 10
        # Python-pptx wraps the XML, so we use xml_slides.insert via the underlying list
        # The element insert is direct on the lxml object
        for i, entry in enumerate(new_entries):
            xml_slides.insert(target_insert_position + i, entry)

    # --- 4. Last slide ("Summary asks — this week") — refresh status ---
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text
                if "Summary asks — this week" in txt or "Five decisions, ranked by ROI" in txt:
                    # Update each decision status
                    for para in shape.text_frame.paragraphs:
                        ptxt = para.text
                        if "Apply LimeSpot 4 admin tweaks" in ptxt and "fw8" in ptxt:
                            new_txt = "Apply LimeSpot 4 admin tweaks — IN PROGRESS [fw8]"
                            if para.runs:
                                para.runs[0].text = new_txt
                                for r in para.runs[1:]: r.text = ""
                        elif "Approve Vertex cutover plan" in ptxt and "oyw" in ptxt:
                            new_txt = "Approve Vertex cutover plan in principle — PENDING [oyw]"
                            if para.runs:
                                para.runs[0].text = new_txt
                                for r in para.runs[1:]: r.text = ""

    # --- 5. Update the last (Questions) slide footer date ---
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Updated May 6, 2026" in run.text:
                            run.text = run.text.replace("Updated May 6, 2026", "Updated May 14, 2026")

    pres.save(str(DECK))
    print(f"Saved {DECK.relative_to(REPO)}")


def main():
    if "--roadmap-only" in sys.argv:
        update_roadmap()
    elif "--deck-only" in sys.argv:
        update_deck()
    else:
        update_roadmap()
        update_deck()
    print("\nDone. Review with:")
    print(f"  git diff --stat {ROADMAP.relative_to(REPO)} {DECK.relative_to(REPO)}")
    print("  open docs/HairMNL-90-Day-Modernization-Roadmap.docx")
    print("  open docs/HairMNL-Theme-Modernization-Deck.pptx")


if __name__ == "__main__":
    sys.exit(main())
